"""Generate a 12-slide Intellidog highlight deck using python-pptx."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent.parent
ASSETS = ROOT / "assets" / "screenshots"
DOCS = ROOT / "docs"
OUT = ROOT / "docs" / "intellidog.pptx"

# Palette
BG = RGBColor(0x0D, 0x11, 0x17)       # github dark
CARD = RGBColor(0x16, 0x1B, 0x22)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)   # blue
GREEN = RGBColor(0x3F, 0xB9, 0x50)
ORANGE = RGBColor(0xE7, 0x6F, 0x51)
RED = RGBColor(0xC1, 0x12, 0x1F)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x94, 0x9E)

W = Inches(13.33)
H = Inches(7.5)


def rgb(r: RGBColor):
    return r


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, left, top, width, height, text, size=Pt(18), bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def heading(slide, title: str, subtitle: str = ""):
    text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             title, size=Pt(36), bold=True, color=ACCENT)
    if subtitle:
        text_box(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
                 subtitle, size=Pt(16), color=MUTED)


def divider(slide, top=Inches(1.25)):
    line = slide.shapes.add_connector(1, Inches(0.6), top, Inches(12.7), top)
    line.line.color.rgb = ACCENT
    line.line.width = Pt(0.75)


def bullet_box(slide, items: list[str], left, top, width, height,
               size=Pt(15), color=WHITE, marker="  -  "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = marker + item
        run.font.size = size
        run.font.color.rgb = color


def image(slide, path: Path, left, top, width, height=None):
    if not path.exists():
        return
    if height:
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(path), left, top, width)


def label(slide, text: str, left, top, width, height, bg_color=CARD, fg_color=WHITE, size=Pt(13)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = fg_color


def build():
    prs = new_prs()

    # ------------------------------------------------------------------
    # Slide 1: Title
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    text_box(s, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
             "Intellidog", size=Pt(56), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1.5), Inches(3.35), Inches(10), Inches(0.6),
             "Intelligent Observability & Event Watchdog", size=Pt(22),
             color=WHITE, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
             "FastAPI  |  Redis  |  SQLite  |  Grafana  |  Claude LLM",
             size=Pt(15), color=MUTED, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1.5), Inches(6.8), Inches(10), Inches(0.4),
             "github.com/rondomondo/intellidog", size=Pt(13),
             color=MUTED, align=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # Slide 2: The Problem
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "The Problem with Static Rules")
    divider(s)
    text_box(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.7),
             "Threshold rules are predictable -- but blind to anything you didn't think to write a rule for.",
             size=Pt(18), color=WHITE)
    bullet_box(s, [
        "error_rate >= 80%  =>  fire alert.  That's it.",
        "No awareness of which service caused it, or what propagated from it.",
        "A DB pool exhaustion causing JWT failures on unrelated services looks like two separate incidents.",
        "On-call engineers spend minutes stitching logs together to find the causal chain.",
        "The more services, the worse the signal-to-noise ratio.",
    ], Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.0), size=Pt(16), color=WHITE)

    # ------------------------------------------------------------------
    # Slide 3: The Core Thesis
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "The Core Thesis: LLM Anomaly Detection")
    divider(s)
    text_box(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
             "Two layers work together. Rules catch what you know. Claude catches what you don't.",
             size=Pt(17), color=MUTED)

    # Left card: rules
    label(s, "Static Alert Rules\n\nFast, deterministic\nThreshold / rate conditions\nYAML-configured\nFires on known patterns",
          Inches(0.6), Inches(2.1), Inches(5.5), Inches(2.8), bg_color=CARD, fg_color=WHITE, size=Pt(15))
    # Right card: LLM
    label(s, "Claude LLM Analyser\n\nRolling event window -> Claude API\nCross-service correlation\nCausal chain reasoning\nFinds what rules miss",
          Inches(7.0), Inches(2.1), Inches(5.7), Inches(2.8), bg_color=CARD, fg_color=ORANGE, size=Pt(15))
    text_box(s, Inches(5.9), Inches(3.1), Inches(1.2), Inches(0.6),
             "vs", size=Pt(30), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    text_box(s, Inches(0.6), Inches(5.1), Inches(12), Inches(0.5),
             "Both layers write alerts to SQLite. Both surface in Grafana. One API key separates them.",
             size=Pt(15), color=MUTED)

    # ------------------------------------------------------------------
    # Slide 4: Mock vs Real LLM
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "MockLLMAnalyser vs Real Claude", "Zero-config convenience vs production insight")
    divider(s)

    text_box(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.4),
             "MockLLMAnalyser (no API key)", size=Pt(15), bold=True, color=GREEN)
    bullet_box(s, [
        "High error rate (>40% critical/high)",
        "Latency spike (duration_ms > 3000ms)",
        "Source burst (>5 events from one service)",
        "Random ambient anomaly on quiet traffic",
    ], Inches(0.6), Inches(1.9), Inches(5.8), Inches(2.5), size=Pt(14), color=WHITE)

    text_box(s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
             "Real Claude -- first live run findings", size=Pt(15), bold=True, color=ORANGE)
    findings = [
        "[critical]  Simultaneous error storm across ALL services",
        "[critical]  DB connection pool exhausted",
        "[high]      Payment timeouts propagating beyond payment-svc",
        "[high]      JWT failures on non-auth services",
        "[medium]    Latency inconsistent with reported errors",
    ]
    bullet_box(s, findings, Inches(7.0), Inches(1.9), Inches(5.8), Inches(3.0),
               size=Pt(13), color=WHITE, marker="")

    text_box(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.7),
             "The mock flags symptoms. Claude identifies the causal chain: "
             "DB exhaustion -> payment timeouts -> JWT failures on unrelated services -- "
             "a cross-service cascade no static rule could detect.",
             size=Pt(13), color=MUTED)

    # ------------------------------------------------------------------
    # Slide 5: Architecture
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "Architecture")
    divider(s)
    arch = DOCS / "architecture.png-1.png"
    if arch.exists():
        image(s, arch, Inches(0.8), Inches(1.35), Inches(11.7))

    # ------------------------------------------------------------------
    # Slide 6: Components
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "Components")
    divider(s)

    col1 = [
        "FastAPI -- ingest, query, alerts, metrics, webhook",
        "Redis pub/sub -- decoupled event fan-out",
        "Alert Engine -- YAML threshold/rate rules",
        "LLM Analyser -- Claude API / MockLLMAnalyser",
    ]
    col2 = [
        "SQLite WAL -- single-file, WAL mode, 30d compaction",
        "Grafana -- SQLite datasource + Infinity datasource",
        "Provisioned alert rules -- error rate, LLM burst",
        "Webhook receiver -- Grafana -> /webhook/grafana loop",
    ]
    bullet_box(s, col1, Inches(0.6), Inches(1.4), Inches(6.0), Inches(3.5), size=Pt(16))
    bullet_box(s, col2, Inches(6.8), Inches(1.4), Inches(6.0), Inches(3.5), size=Pt(16))

    # ------------------------------------------------------------------
    # Slide 7: Quick Start
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "Quick Start", "Three commands to a running system")
    divider(s)

    code = [
        "git clone git@github.com:rondomondo/intellidog.git && cd intellidog",
        "",
        "make up          # builds images, starts API :8000 + Grafana :3000 + Redis",
        "",
        "make install     # install local dev tools (generate/query scripts)",
        "make generate    # POST 20 synthetic events",
        "make generate-spike   # POST 60 critical events -- triggers alert rules",
        "",
        "open http://localhost:8000   # mission control index",
        "open http://localhost:3000   # Grafana  (admin / admin)",
    ]
    text_box(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5),
             "\n".join(code), size=Pt(15), color=GREEN)

    # ------------------------------------------------------------------
    # Slide 8: Mission Control (index page screenshot)
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "Mission Control: localhost:8000")
    divider(s)
    text_box(s, Inches(0.6), Inches(1.35), Inches(5.2), Inches(0.6),
             "The API root is a dark-themed HTML index -- the first stop after make up. "
             "If it loads, the API and DB are up. Navigate directly to /health, /docs, "
             "/metrics/summary, or any query endpoint.",
             size=Pt(14), color=MUTED)
    idx = ASSETS / "index_page.png"
    image(s, idx, Inches(6.0), Inches(1.35), Inches(6.8))

    # ------------------------------------------------------------------
    # Slide 9: Grafana Dashboard
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "Grafana Dashboard", "Fully provisioned -- no manual setup")
    divider(s)
    dash = ASSETS / "grafana_dashboard.png"
    image(s, dash, Inches(0.4), Inches(1.35), Inches(8.5))
    bullet_box(s, [
        "Event Rate over time",
        "Error Rate %",
        "Severity breakdown",
        "P95 latency",
        "Recent alerts table",
        "LLM anomaly count",
    ], Inches(9.1), Inches(1.8), Inches(4.0), Inches(4.0), size=Pt(14))

    # ------------------------------------------------------------------
    # Slide 10: Alert Pipeline
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "End-to-End Alert Pipeline")
    divider(s)
    bullet_box(s, [
        "Events ingested via POST /events -- written to SQLite, published to Redis",
        "Alert Engine evaluates YAML rules on every event (threshold, rate, window)",
        "LLM Analyser runs every 60s -- sends rolling window to Claude, stores findings",
        "Grafana queries /metrics/summary via Infinity datasource (server-side eval)",
        "Grafana fires alert -> POST /webhook/grafana -> stored, visible at GET /webhook/grafana",
        "Provisioned rules: error_rate_pct > 80%  |  llm_anomalies_last_hour > 10",
    ], Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.0), size=Pt(16))

    alerts = ASSETS / "grafana_alerts_panel.png"
    image(s, alerts, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.3))

    # ------------------------------------------------------------------
    # Slide 11: Testing
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    heading(s, "Testing", "92% coverage, 116 tests, zero external dependencies")
    divider(s)
    bullet_box(s, [
        "pytest-asyncio + httpx.AsyncClient -- full in-process FastAPI testing",
        "All Redis, SQLite, and Anthropic API calls mocked in tests",
        "MockLLMAnalyser exercises the full alert pipeline in CI without an API key",
        "make test  |  make ci  (format + lint + typecheck + test)",
    ], Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.2), size=Pt(16))

    modules = [
        ("api/health.py", "100%"), ("api/ingest.py", "100%"),
        ("bus/publisher.py", "100%"), ("db/connection.py", "100%"),
        ("api/metrics.py", "97%"), ("api/query.py", "94%"),
        ("bus/subscriber.py", "93%"), ("engine/alert_engine.py", "91%"),
        ("engine/llm_analyser.py", "88%"), ("db/repository.py", "89%"),
    ]
    col_w = Inches(3.0)
    row_h = Inches(0.42)
    for i, (mod, cov) in enumerate(modules):
        col = i % 2
        row = i // 2
        left = Inches(0.6) + col * Inches(6.5)
        top = Inches(3.7) + row * row_h
        label(s, f"{mod}  --  {cov}", left, top, Inches(5.8), row_h - Inches(0.04),
              bg_color=CARD, fg_color=GREEN if cov == "100%" else WHITE, size=Pt(13))

    # ------------------------------------------------------------------
    # Slide 12: Summary
    # ------------------------------------------------------------------
    s = blank_slide(prs)
    bg(s)
    text_box(s, Inches(1.0), Inches(1.6), Inches(11), Inches(1.0),
             "Intellidog", size=Pt(44), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    text_box(s, Inches(1.0), Inches(2.55), Inches(11), Inches(0.5),
             "Self-hosted observability with real LLM anomaly intelligence", size=Pt(19),
             color=WHITE, align=PP_ALIGN.CENTER)

    bullet_box(s, [
        "Drop-in event ingestion -- any service, any language, one POST",
        "YAML alert rules -- no code, no redeploy",
        "MockLLMAnalyser works out of the box -- add a key to go real",
        "Full Grafana integration -- dashboard, alerts, and webhook loop provisioned",
        "92% test coverage -- CI-ready from day one",
    ], Inches(2.0), Inches(3.3), Inches(9.5), Inches(2.8), size=Pt(17))

    text_box(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.5),
             "make up  |  make generate  |  open http://localhost:8000",
             size=Pt(15), color=GREEN, align=PP_ALIGN.CENTER)

    prs.save(str(OUT))
    print(f"Saved {OUT}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
